"""
ingest.py — Standalone script to ingest business-law PDF and PPTX materials into LightRAG.

Usage:
    python ingest.py

Reads all PDFs, PPTX, and DOCX files from ../data/materials/business-law/ using pymupdf (fitz),
python-pptx, and python-docx, inserts text into LightRAG (lightrag-hku >= 1.3.8) using OpenAI
gpt-4o-mini + text-embedding-3-small.
"""

import os
import asyncio
import fitz  # pymupdf
from pptx import Presentation
from docx import Document

from dotenv import load_dotenv
from lightrag import LightRAG
from lightrag.llm.openai import gpt_4o_mini_complete, openai_embed
from lightrag.utils import EmbeddingFunc

load_dotenv()

# Resolve paths relative to this script's location
BACKEND_DIR = os.path.dirname(os.path.abspath(__file__))
DATA_DIR = os.path.join(BACKEND_DIR, "..", "data")
MATERIALS_DIR = os.path.join(DATA_DIR, "materials", "business-law")
LIGHTRAG_DIR = os.path.join(DATA_DIR, "lightrag_storage")

os.makedirs(LIGHTRAG_DIR, exist_ok=True)


def extract_text_from_pdf(pdf_path: str) -> str:
    """Extract all text from a PDF file using pymupdf."""
    doc = fitz.open(pdf_path)
    pages = []
    for page in doc:
        pages.append(page.get_text())
    doc.close()
    return "\n".join(pages)


def extract_text_from_pptx(path: str) -> str:
    """Extract all text from a PPTX file using python-pptx."""
    prs = Presentation(path)
    slides_text = []
    for slide in prs.slides:
        shapes_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                shapes_text.append(shape.text)
        slides_text.append("\n".join(shapes_text))
    return "\n".join(slides_text)


def extract_text_from_docx(path: str) -> str:
    """Extract all text from a DOCX file using python-docx."""
    doc = Document(path)
    paragraphs_text = []
    for paragraph in doc.paragraphs:
        paragraphs_text.append(paragraph.text)
    return "\n".join(paragraphs_text)


async def ingest_all():
    print(f"LightRAG working dir: {os.path.abspath(LIGHTRAG_DIR)}")
    print(f"Materials dir: {os.path.abspath(MATERIALS_DIR)}")

    rag = LightRAG(
        working_dir=LIGHTRAG_DIR,
        llm_model_func=gpt_4o_mini_complete,
        embedding_func=EmbeddingFunc(
            embedding_dim=1536,
            max_token_size=8192,
            func=lambda texts: openai_embed(texts, model="text-embedding-3-small"),
        ),
    )
    await rag.initialize_storages()

    all_files = sorted([
        f for f in os.listdir(MATERIALS_DIR)
        if f.lower().endswith(".pdf") or f.lower().endswith(".pptx") or f.lower().endswith(".docx")
    ])

    if not all_files:
        print(f"No PDF, PPTX, or DOCX files found in {MATERIALS_DIR}")
        return

    print(f"Found {len(all_files)} file(s): {all_files}\n")

    for filename in all_files:
        file_path = os.path.join(MATERIALS_DIR, filename)
        print(f"[>>] Processing: {filename}")
        try:
            if filename.lower().endswith(".pdf"):
                text = extract_text_from_pdf(file_path)
            elif filename.lower().endswith(".pptx"):
                text = extract_text_from_pptx(file_path)
            else:
                text = extract_text_from_docx(file_path)
            if not text.strip():
                print(f"  [!] Skipped (empty text): {filename}")
                continue
            print(f"  [i] Extracted {len(text)} characters")
            await rag.ainsert(text)
            print(f"  [OK] Inserted into LightRAG: {filename}")
        except Exception as e:
            print(f"  [ERR] Error processing {filename}: {e}")

    print("\nIngestion complete.")


if __name__ == "__main__":
    asyncio.run(ingest_all())
