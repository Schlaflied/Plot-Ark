"""File text extraction: PDF, PPTX, DOCX."""

import os
import io
import re
import fitz  # pymupdf
import docx as _docx_lib


def detect_module_from_pptx(filename: str, content: bytes) -> int | None:
    """Return module number inferred from filename or first-slide title, or None."""
    # Layer 1: filename pattern — M1, Module1, Week1, Lecture1, L1, etc.
    name = os.path.splitext(os.path.basename(filename))[0]
    m = re.search(r'(?:module|mod|week|lecture|lec|unit|m|w|l)[\s_\-]?(\d+)', name, re.IGNORECASE)
    if m:
        return int(m.group(1))

    # Layer 2: first slide title shape
    try:
        from pptx import Presentation
        from pptx.util import Pt
        prs = Presentation(io.BytesIO(content))
        if prs.slides:
            slide = prs.slides[0]
            for shape in slide.shapes:
                if shape.has_text_frame and shape.shape_type == 13:  # title placeholder
                    text = shape.text
                    m = re.search(r'(?:module|mod|week|lecture|lec|unit)[\s_\-]?(\d+)', text, re.IGNORECASE)
                    if m:
                        return int(m.group(1))
            # fallback: any shape on first slide
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text
                    m = re.search(r'(?:module|mod|week|lecture|lec|unit)[\s_\-]?(\d+)', text, re.IGNORECASE)
                    if m:
                        return int(m.group(1))
    except Exception:
        pass

    return None


def extract_text_from_bytes(filename: str, content: bytes) -> str:
    """Extract plain text from PDF, PPTX, or DOCX bytes."""
    ext = os.path.splitext(filename.lower())[1]
    if ext == ".pdf":
        doc = fitz.open(stream=content, filetype="pdf")
        pages = [page.get_text() for page in doc]
        doc.close()
        return "\n".join(pages)
    elif ext == ".pptx":
        from pptx import Presentation
        prs = Presentation(io.BytesIO(content))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    texts.append(shape.text)
        return "\n".join(texts)
    elif ext == ".docx":
        doc = _docx_lib.Document(io.BytesIO(content))
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    return ""
