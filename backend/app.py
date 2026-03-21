import os
import re
import json
import time
import uuid
import asyncio
import tempfile
import psycopg2
from flask import Flask, request, Response, stream_with_context, jsonify, send_file
import fitz  # pymupdf
import docx as _docx_lib
import io
from flask_cors import CORS
from dotenv import load_dotenv
from openai import OpenAI
import google.generativeai as genai
from tavily import TavilyClient

# ---------------------------------------------------------------------------
# Module-level caches
# ---------------------------------------------------------------------------
_rag_instances = {}       # key: storage_dir path → LightRAG instance
_initialized_instances = set()  # storage_dirs that have had initialize_storages() called
_ingest_jobs = {}         # key: job_id → {"status": "running"|"done"|"error", "progress": str, "message": str}

# Persistent background event loop — never closed, so LightRAG's internal state stays valid
import threading as _threading
_bg_loop = asyncio.new_event_loop()
_bg_thread = _threading.Thread(target=_bg_loop.run_forever, daemon=True)
_bg_thread.start()

def _run_async(coro):
    """Submit a coroutine to the persistent background event loop and wait for result."""
    future = asyncio.run_coroutine_threadsafe(coro, _bg_loop)
    return future.result(timeout=120)

try:
    import redis as _redis_lib
    _redis_client = _redis_lib.Redis(host="redis", port=6379, db=0, decode_responses=True)
    _redis_client.ping()
    print("Redis cache connected.")
except Exception as _redis_err:
    print(f"Redis unavailable, caching disabled: {_redis_err}")
    _redis_client = None

load_dotenv()

app = Flask(__name__)
CORS(app, resources={r"/api/*": {"origins": "*"}})

AI_PROVIDER = os.getenv("AI_PROVIDER", "openai").lower()
openai_client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))
genai.configure(api_key=os.getenv("GEMINI_API_KEY", ""))
tavily_client = TavilyClient(api_key=os.getenv("TAVILY_API_KEY"))

DATABASE_URL = os.getenv("DATABASE_URL", "postgresql://plotark:plotark@postgres:5432/plotark")


def get_db():
    try:
        return psycopg2.connect(DATABASE_URL)
    except Exception as e:
        print(f"DB connection error: {e}")
        return None


def init_db():
    for attempt in range(10):
        conn = get_db()
        if conn:
            try:
                cur = conn.cursor()
                cur.execute("""
                    CREATE TABLE IF NOT EXISTS curricula (
                        id SERIAL PRIMARY KEY,
                        created_at TIMESTAMP DEFAULT NOW(),
                        topic TEXT NOT NULL,
                        level TEXT,
                        audience TEXT,
                        course_code TEXT,
                        course_type TEXT,
                        module_count INTEGER,
                        modules JSONB,
                        sources JSONB,
                        is_favorite BOOLEAN DEFAULT FALSE
                    )
                """)
                cur.execute("""
                    ALTER TABLE curricula ADD COLUMN IF NOT EXISTS is_favorite BOOLEAN DEFAULT FALSE
                """)
                cur.execute("""
                    CREATE TABLE IF NOT EXISTS xapi_statements (
                        id SERIAL PRIMARY KEY,
                        actor_email TEXT NOT NULL,
                        actor_name TEXT NOT NULL,
                        verb TEXT NOT NULL,
                        object_id TEXT NOT NULL,
                        object_name TEXT NOT NULL,
                        timestamp TIMESTAMPTZ NOT NULL DEFAULT NOW(),
                        curriculum_topic TEXT
                    )
                """)
                conn.commit()
                cur.close()
                conn.close()
                print("DB initialized.")
                return
            except Exception as e:
                print(f"DB init error: {e}")
                conn.close()
                return
        print(f"DB not ready, retrying ({attempt + 1}/10)...")
        time.sleep(3)
    print("Could not connect to DB after 10 attempts. Continuing without DB.")


def save_curriculum(topic, level, audience, course_code, course_type, module_count, data, design_approach="addie"):
    conn = get_db()
    if not conn:
        return
    try:
        cur = conn.cursor()
        cur.execute(
            """INSERT INTO curricula (topic, level, audience, course_code, course_type, module_count, modules, sources)
               VALUES (%s, %s, %s, %s, %s, %s, %s, %s)""",
            (topic, level, audience, course_code, course_type, module_count,
             json.dumps(data.get("modules", [])),
             json.dumps(data.get("sources", [])))
        )
        conn.commit()
        cur.close()
        conn.close()
    except Exception as e:
        print(f"DB save error: {e}")


BLOOMS_BY_LEARNER_LEVEL = {
    "beginner": {
        "label": "Remember and Understand",
        "verbs": "define, identify, recall, describe, explain, summarize",
        "constraint": "Learning objectives MUST use only Remember and Understand verbs: define, identify, recall, describe, explain, summarize. Do NOT use Apply, Analyze, Evaluate, or Create verbs.",
    },
    "intermediate": {
        "label": "Apply and Analyze",
        "verbs": "apply, demonstrate, differentiate, compare, examine, solve",
        "constraint": "Learning objectives MUST use only Apply and Analyze verbs: apply, demonstrate, differentiate, compare, examine, solve. Do NOT use Remember, Understand, Evaluate, or Create verbs.",
    },
    "advanced": {
        "label": "Evaluate and Create",
        "verbs": "assess, critique, design, construct, argue, justify, synthesize",
        "constraint": "Learning objectives MUST use only Evaluate and Create verbs: assess, critique, design, construct, argue, justify, synthesize. Do NOT use Remember, Understand, Apply, or Analyze verbs.",
    },
}


LEVEL_TO_BLOOMS = {
    # Undergraduate
    "undergraduate-year-1": "beginner",
    "undergraduate-year-2": "beginner",
    "undergraduate-year-3": "intermediate",
    "undergraduate-year-4": "intermediate",
    # Graduate
    "master-year-1": "advanced",
    "master-year-2": "advanced",
    "master-year-3": "advanced",
    "doctoral": "advanced",
    # Professional
    "professional-beginner": "beginner",
    "professional-intermediate": "intermediate",
    "professional-advanced": "advanced",
    # ESL/EFL
    "esl-beginner": "beginner",
    "esl-intermediate": "intermediate",
    "esl-advanced": "advanced",
    # K-12
    "k12-elementary": "beginner",
    "k12-middle": "beginner",
    "k12-highschool": "intermediate",
}

_BLOOMS_TO_NARRATIVE = {
    "beginner": "Remember and Understand — definitions, identification, basic comprehension",
    "intermediate": "Apply and Analyze — case analysis, pattern recognition, comparative evaluation",
    "advanced": "Analyze, Evaluate, and Create — synthesis, critique, independent judgment, original work",
}


def get_blooms_level(course_code, level):
    # Check structured level key first
    if str(level) in LEVEL_TO_BLOOMS:
        return _BLOOMS_TO_NARRATIVE[LEVEL_TO_BLOOMS[str(level)]]
    num_match = re.search(r'\d{3}', str(course_code).upper())
    if num_match:
        num = int(num_match.group())
        if num < 200:
            return "Remember and Understand — definitions, identification, basic comprehension"
        elif num < 300:
            return "Understand and Apply — using concepts in familiar contexts, worked examples"
        elif num < 400:
            return "Apply and Analyze — case analysis, pattern recognition, comparative evaluation"
        else:
            return "Analyze, Evaluate, and Create — synthesis, critique, independent judgment, original work"
    level_lower = str(level).lower()
    if any(k in level_lower for k in ['graduate', 'phd', 'doctoral', 'master']):
        return "Analyze, Evaluate, and Create — advanced critical synthesis and original contribution"
    elif any(k in level_lower for k in ['senior', '4th', 'advanced']):
        return "Apply, Analyze, and Evaluate — critical analysis with some synthesis"
    return "Understand and Apply — foundational understanding with practical application"


def get_session_constraints(minutes):
    """Return a prompt instruction string based on session duration in minutes."""
    if minutes <= 75:
        return (
            "Session length: 75 minutes. "
            "Each module must be completable in 75 minutes. "
            "Max 1 required reading per module (≤15 min read). "
            "Assignments must be short and focused (≤30 min completion time). "
            "Prefer in-class discussion or quick reflection over lengthy projects."
        )
    elif minutes <= 90:
        return (
            "Session length: 90 minutes. "
            "Each module fits a standard 90-minute university class. "
            "Max 1-2 required readings per module (≤20 min read each). "
            "Assignments should be completable in 45-60 minutes."
        )
    else:  # 3 hours or more
        return (
            f"Session length: {minutes} minutes (extended format). "
            "Each module covers more ground with deeper engagement. "
            "Up to 2-3 readings allowed per module. "
            "Assignments can include workshop components, group activities, or multi-part tasks. "
            "Include at least one in-class activity suggestion per module."
        )


def get_blooms_constraint(level):
    """Return Bloom's verb constraint based on beginner/intermediate/advanced learner level."""
    # Check structured level key first
    if str(level) in LEVEL_TO_BLOOMS:
        return BLOOMS_BY_LEARNER_LEVEL[LEVEL_TO_BLOOMS[str(level)]]["constraint"]
    level_lower = str(level).lower()
    if level_lower in BLOOMS_BY_LEARNER_LEVEL:
        return BLOOMS_BY_LEARNER_LEVEL[level_lower]["constraint"]
    # Fuzzy fallback
    if any(k in level_lower for k in ['begin', 'intro', 'foundation', '100', '1st', 'first']):
        return BLOOMS_BY_LEARNER_LEVEL["beginner"]["constraint"]
    if any(k in level_lower for k in ['advanc', 'senior', 'graduate', 'expert', 'master', 'phd', 'doctoral']):
        return BLOOMS_BY_LEARNER_LEVEL["advanced"]["constraint"]
    return BLOOMS_BY_LEARNER_LEVEL["intermediate"]["constraint"]


RESOURCE_TYPES = {
    "academic": {
        "domains": ["jstor.org", "researchgate.net", "academia.edu", "ncbi.nlm.nih.gov",
                    "springer.com", "tandfonline.com", "sagepub.com", "wiley.com",
                    "oxfordhandbooks.com", "cambridge.org", "scholar.google.com"],
        "queries": [
            "{topic} academic research {level}",
            "{topic} {audience} course materials",
            "{topic} key concepts textbook",
        ],
        "max_per_query": 3,
    },
    "video": {
        "domains": ["youtube.com", "ted.com", "coursera.org", "edx.org", "khanacademy.org"],
        "queries": [
            "{topic} lecture video course",
            "{topic} TED talk introduction",
        ],
        "max_per_query": 2,
    },
    "news": {
        "domains": ["hbr.org", "economist.com", "nytimes.com", "theguardian.com",
                    "mit.edu", "stanford.edu", "bbc.com"],
        "queries": [
            "{topic} analysis report {level}",
        ],
        "max_per_query": 2,
    },
}


def research_sources(topic, level, audience):
    """Step 1: Agent searches for real sources by type before generation."""
    try:
        results = []
        for source_type, config in RESOURCE_TYPES.items():
            for query_template in config["queries"]:
                query = query_template.format(topic=topic, level=level, audience=audience)
                search_kwargs = {
                    "query": query,
                    "search_depth": "basic",
                    "max_results": config["max_per_query"],
                    "include_domains": config["domains"],
                }
                if source_type == "news":
                    search_kwargs["days"] = 365
                response = tavily_client.search(**search_kwargs)
                for r in response.get("results", []):
                    results.append({
                        "title": r.get("title", ""),
                        "url": r.get("url", ""),
                        "content": r.get("content", "")[:300],
                        "type": source_type,
                    })
        # Deduplicate by URL
        seen = set()
        unique = []
        for r in results:
            if r["url"] not in seen:
                seen.add(r["url"])
                unique.append(r)
        # Filter: drop sources with no title or content clearly unrelated to topic
        topic_keywords = set(topic.lower().split())
        def is_relevant(r):
            if not r["title"]:
                return False
            combined = (r["title"] + " " + r["content"]).lower()
            return any(kw in combined for kw in topic_keywords)
        filtered = [r for r in unique if is_relevant(r)]
        if len(filtered) < 3:
            filtered = unique  # fallback: keep all if filter too aggressive
        print(f"Tavily found {len(unique)} sources, {len(filtered)} passed relevance filter for: {topic}")
        return filtered[:10]
    except Exception as e:
        print(f"Tavily research error: {e}")
        return []


ACADEMIC_DOMAINS = [
    "jstor.org", "springer.com", "researchgate.net", "sciencedirect.com",
    "wiley.com", "pubmed.ncbi.nlm.nih.gov", "scholar.google.com", ".edu",
]
NEWS_DOMAINS = [
    "nytimes.com", "economist.com", "hbr.org", "theguardian.com",
    "bbc.com", "reuters.com", ".gov", ".gc.ca",
]
VIDEO_DOMAINS = [
    "youtube.com", "ted.com", "coursera.org", "edx.org",
]


def score_credibility(url: str, source_type: str) -> str:
    """Return 'high', 'medium', or 'low' credibility based on domain heuristics."""
    url_lower = url.lower()
    if source_type == "academic":
        if any(d in url_lower for d in ACADEMIC_DOMAINS):
            return "high"
    if source_type == "news":
        if any(d in url_lower for d in NEWS_DOMAINS):
            return "medium"
    if source_type == "video":
        if any(d in url_lower for d in VIDEO_DOMAINS):
            return "medium"
    # Cross-check all domains regardless of source_type label
    if any(d in url_lower for d in ACADEMIC_DOMAINS):
        return "high"
    if any(d in url_lower for d in NEWS_DOMAINS + VIDEO_DOMAINS):
        return "medium"
    return "low"


ASSESSMENT_FORMATS = {
    "project": "project-based assignments (group projects, case studies, presentations, portfolios)",
    "essay": "essay-based assessments (argumentative essays, reflective journals, research papers)",
    "debate": "discussion-based formats (structured debates, Socratic seminars, roleplay scenarios)",
    "lab": "lab and simulation work (experiments, technical projects, lab reports, prototypes)",
    "mixed": "varied formats across modules (rotate between essays, projects, discussions, and activities)",
}


@app.route("/")
def index():
    return {"status": "online", "service": "Plot Ark — Agentic Curriculum Engine"}


@app.route("/api/sources/preview", methods=["POST"])
def preview_sources():
    """Return Tavily sources for user review before curriculum generation."""
    data = request.get_json()
    topic = data.get("topic", "")
    level = data.get("level", "")
    audience = data.get("audience", "")

    if not all([topic, level, audience]):
        return {"error": "Missing required fields: topic, level, audience"}, 400

    # --- Redis cache check ---
    cache_key = f"sources_preview:{topic}:{level}:{audience}"
    if _redis_client is not None:
        try:
            cached = _redis_client.get(cache_key)
            if cached:
                print(f"Redis cache hit: {cache_key}")
                return json.loads(cached)
        except Exception as redis_err:
            print(f"Redis get error (skipping cache check): {redis_err}")

    raw = research_sources(topic, level, audience)
    sources = []
    for r in raw:
        sources.append({
            "url": r.get("url", ""),
            "title": r.get("title", ""),
            "type": r.get("type", "other"),
            "snippet": r.get("content", ""),
            "credibility": score_credibility(r.get("url", ""), r.get("type", "")),
            "tags": [],
        })

    # Batch-generate keyword tags for all sources in one GPT call
    if sources:
        try:
            sources_for_tags = [{"title": s["title"], "snippet": s["snippet"]} for s in sources]
            tag_prompt = (
                "For each of the following academic sources, generate:\n"
                "1. 3-4 short keyword tags (1-3 words each)\n"
                "2. A clean one-sentence summary (max 20 words) describing what the source covers — no filler like 'This paper examines', just the actual content\n\n"
                "Return as JSON only, no explanation.\n\nSources:\n"
                + json.dumps(sources_for_tags)
                + '\n\nReturn format: {"results": [{"tags": ["tag1", "tag2"], "summary": "one sentence"}, ...]}\n'
                "Each object corresponds to one source in the same order."
            )

            tag_response = openai_client.chat.completions.create(
                model="gpt-4o-mini",
                messages=[{"role": "user", "content": tag_prompt}],
                temperature=0.3,
                max_tokens=500,
            )
            raw_tags = tag_response.choices[0].message.content.strip()
            # Strip markdown code fences if present
            raw_tags = raw_tags.replace("```json", "").replace("```", "").strip()
            tag_data = json.loads(raw_tags)
            results_list = tag_data.get("results", [])
            for i, source in enumerate(sources):
                if i < len(results_list) and isinstance(results_list[i], dict):
                    source["tags"] = results_list[i].get("tags", [])
                    source["snippet"] = results_list[i].get("summary", source["snippet"])
        except Exception as tag_err:
            print(f"Tag generation failed (non-fatal): {tag_err}")
            # Fall back to empty tags — already set above

    # --- Store result in Redis ---
    if _redis_client is not None:
        try:
            _redis_client.setex(cache_key, 604800, json.dumps({"sources": sources}))  # 7 days
        except Exception as redis_err:
            print(f"Redis set error (skipping cache store): {redis_err}")

    return {"sources": sources}


@app.route("/api/curriculum/generate", methods=["POST"])
def generate_curriculum():
    data = request.get_json()
    topic = data.get("topic", "")
    level = data.get("level", "")
    audience = data.get("audience", "")
    accreditation_context = data.get("accreditation_context", "")
    course_code = data.get("course_code", "")
    course_type = data.get("course_type", "mixed")
    module_count_raw = data.get("module_count", "6")
    design_approach = data.get("design_approach", "addie").lower()
    if design_approach not in ("addie", "sam"):
        design_approach = "addie"

    try:
        session_duration = max(1, int(data.get("session_duration", 90)))
    except (ValueError, TypeError):
        session_duration = 90

    if not all([topic, level, audience]):
        return {"error": "Missing required fields"}, 400

    try:
        module_count = max(3, min(12, int(module_count_raw)))
    except (ValueError, TypeError):
        module_count = 6

    blooms = get_blooms_level(course_code, level)
    blooms_constraint = get_blooms_constraint(level)
    session_constraint = get_session_constraints(session_duration)
    assessment_format = ASSESSMENT_FORMATS.get(course_type, ASSESSMENT_FORMATS["mixed"])

    # Step 1: Use approved_sources if provided (R2 human-in-the-loop), otherwise run Tavily
    approved_sources_raw = data.get("approved_sources", None)
    required_sources = []
    optional_sources = []
    if approved_sources_raw and isinstance(approved_sources_raw, list) and len(approved_sources_raw) > 0:
        # Convert approved_sources format (url/title/type/snippet/priority) to internal format
        real_sources = []
        for s in approved_sources_raw:
            if not s.get("url"):
                continue
            priority = s.get("priority", "optional")
            entry = {
                "url": s.get("url", ""),
                "title": s.get("title", ""),
                "type": s.get("type", "other"),
                "content": s.get("snippet", ""),
                "priority": priority,
            }
            real_sources.append(entry)
            if priority == "required":
                required_sources.append(entry)
            else:
                optional_sources.append(entry)
        print(f"Using {len(real_sources)} user-approved sources ({len(required_sources)} required, {len(optional_sources)} optional) — skipping Tavily")
    else:
        real_sources = research_sources(topic, level, audience)
    sources_context = ""
    if real_sources:
        sources_context = "\n\nReal sources found by research agent — use these URLs in your sources array (they are verified real):\n"
        for s in real_sources:
            priority_label = s.get("priority", "")
            priority_tag = f" [PRIORITY: {priority_label.upper()}]" if priority_label else ""
            sources_context += f"- [{s['type']}]{priority_tag} {s['title']} | {s['url']}\n"
        sources_context += "\nPrioritize these real URLs. You may add more you know with confidence, but do NOT invent URLs.\n"

    # Build reading priority instructions for the prompt
    reading_priority_instructions = ""
    if required_sources or optional_sources:
        reading_priority_instructions = "\n\nReading Priority Instructions (based on instructor selection):\n"
        if required_sources:
            reading_priority_instructions += "REQUIRED readings — these MUST appear in modules as assigned readings:\n"
            for s in required_sources:
                reading_priority_instructions += f"  - {s['title']} | {s['url']}\n"
        if optional_sources:
            reading_priority_instructions += "OPTIONAL/supplementary readings — include where relevant but not mandatory:\n"
            for s in optional_sources:
                reading_priority_instructions += f"  - {s['title']} | {s['url']}\n"
        reading_priority_instructions += (
            "When assigning readings to modules:\n"
            "- Mark required readings with \"reading_type\": \"required\"\n"
            "- Mark optional readings with \"reading_type\": \"optional\"\n"
        )
    else:
        reading_priority_instructions = (
            "\n\nFor each reading in recommended_readings, assign a reading_type field:\n"
            "- \"required\" if it directly covers the core concept of the module\n"
            "- \"optional\" if it is supplementary or extension material\n"
        )

    # Build design-approach-specific instructions
    if design_approach == "sam":
        design_approach_label = "SAM (Successive Approximation Model)"
        design_approach_instructions = """
Design Approach — SAM (Successive Approximation Model):
- Frame each module with ITERATIVE checkpoints rather than fixed deliverables.
- Each module MUST include a "rapid_prototype_cycle" field: a brief description of the Rapid Prototype → Evaluate → Revise loop learners will go through.
- Assignments should be framed as low-stakes prototypes designed to be revised, not final submissions.
- The overall curriculum narrative should emphasize continuous iteration over linear completion.
"""
        sam_module_field = '''"rapid_prototype_cycle": "Description of the Rapid Prototype → Evaluate → Revise cycle for this module.",'''
    else:
        design_approach_label = "ADDIE (Analysis → Design → Development → Implementation → Evaluation)"
        design_approach_instructions = """
Design Approach — ADDIE (linear instructional design model):
- Follow the standard linear flow: Analysis → Design → Development → Implementation → Evaluation.
- Each module represents a discrete, completed stage of learning before the next begins.
- Assignments are summative deliverables that demonstrate mastery of that module's objectives.
"""
        sam_module_field = ""

    resource_priority_prompt = {
        "project": "RESOURCE PRIORITY: Each module's recommended_readings MUST include at least 1 news or industry source (HBR, Economist, NYT, etc.) alongside academic sources. Real-world cases are essential for project-based courses.",
        "essay": "RESOURCE PRIORITY: Prioritize academic sources. Add video (TED Talk, lecture) and news sources where they strengthen the argument. Never omit academic sources.",
        "debate": "RESOURCE PRIORITY: Each module MUST include at least 1 current news or policy source to support debate positions. Mix with academic sources for theoretical grounding.",
        "lab": "RESOURCE PRIORITY: Each module MUST include at least 1 video resource (tutorial, demonstration, simulation walkthrough). Supplement with academic readings.",
        "mixed": "RESOURCE PRIORITY: Distribute resource types across modules — not every module should be academic-only. Mix academic, news (current events), and video across the curriculum.",
    }

    prompt = f"""You are an expert curriculum designer applying evidence-based instructional design principles. Generate a rigorous, narrative-driven curriculum.

Topic: {topic}
Course Code: {course_code or "Not specified"}
Level: {level}
Target Audience: {audience}
Accreditation Context: {accreditation_context}
Course Type: {course_type}
Number of Modules: {module_count}
Design Approach: {design_approach_label}

Pedagogical Constraints:
- Bloom's Taxonomy Target: {blooms}
- Bloom's Verb Constraint: {blooms_constraint}
- Session Duration: {session_constraint}
- Assessment Format: {assessment_format}
- Difficulty Progression (i+1 principle, Krashen): complexity_level must start at 1 and reach 5 by the final module, increasing evenly — never jump more than 1 level per module.
- Cognitive Load (Sweller): Maximum 2 recommended readings per module. Each reading must have a clear rationale tied to that module's learning objectives.
- Not every module requires an assignment. When included, it must align with the module's Bloom's level and course type.
- Assignment task_description: MUST be specific and actionable (e.g. "Write a 500-word reflection comparing two case studies..."), NOT generic (e.g. "This assignment addresses the objectives of..."). Failing this instruction makes the output unusable.
- Assignment rubric_highlights: MUST contain exactly 3-4 concrete criteria describing what excellent work looks like for THIS specific task.
- Assignment estimated_time: MUST be realistic given the session duration constraint above. A 75-min session cannot have a 3-hour assignment.
{design_approach_instructions}{reading_priority_instructions}
{resource_priority_prompt.get(course_type, resource_priority_prompt["mixed"])}
Return ONLY valid JSON (no markdown, no explanation):
{{
  "design_approach": "{design_approach}",
  "session_duration_minutes": {session_duration},
  "modules": [
    {{
      "title": "Module title",
      "complexity_level": 1,
      "learning_objectives": ["objective using only the permitted Bloom's verbs for this level", "objective 2", "objective 3"],
      {sam_module_field}
      "narrative_preview": "A compelling 2-3 sentence narrative hook using metaphor, scenario, or challenge framing.",
      "recommended_readings": [
        {{
          "title": "Full title of reading (article, chapter, or textbook section) — complete, never truncated with '...'",
          "url": "https://real-url-from-sources-above.com",
          "type": "academic | video | news",
          "estimated_time": "15 min read | 20 min video | 10 min read",
          "reading_type": "required | optional",
          "key_points": ["key point 1", "key point 2"],
          "rationale": "Why this reading is essential for this module's specific learning objectives and why it is relevant to students' lives or careers."
        }}
      ],
      "assignments": [
        {{
          "type": "project | essay | quiz | discussion | presentation | lab | reflection",
          "title": "Short assignment title",
          "task_description": "2-3 sentence specific description of exactly what students must do. Must be concrete and actionable — NOT generic phrases like 'addresses the objectives of this module'.",
          "deliverable": "What they hand in — e.g. '1-page written reflection', '10-slide deck', 'in-class oral presentation (5 min)'",
          "estimated_time": "Realistic completion time given the session duration — e.g. '45 minutes', '2 hours'. Must match the session length constraint above.",
          "covers_objectives": "Which specific learning objectives from this module this assessment addresses.",
          "rubric_highlights": [
            "Criterion 1 — description of what good work looks like",
            "Criterion 2 — description of what good work looks like",
            "Criterion 3 — description of what good work looks like",
            "Criterion 4 — description of what good work looks like"
          ]
        }}
      ]
    }}
  ],
  "sources": [
    {{
      "title": "Full title of the paper, video, article, or resource — complete, never truncated with '...'",
      "url": "https://example.com",
      "domain": "example.com",
      "type": "academic | video | news",
      "estimated_time": "20 min read | 15 min video | 10 min read",
      "retrieved_at": "2026-03-16"
    }}
  ]
}}

CRITICAL REQUIREMENT — MODULE COUNT:
You MUST generate exactly {module_count} modules. No more, no fewer.
Before finalizing your response, count your modules and verify the count equals {module_count}.
If your count is wrong, fix it before responding.
Responses with incorrect module counts will be automatically rejected and regenerated.

complexity_level must start at 1 and reach 5 by the last module.

For sources: use the verified real URLs provided above. Add more real sources you know with confidence. Every URL must be real and accessible.{sources_context}"""

    def event_stream():
        if approved_sources_raw and isinstance(approved_sources_raw, list) and len(approved_sources_raw) > 0:
            yield f"data: {json.dumps({'status': 'generating', 'message': f'Generating curriculum with {len(approved_sources_raw)} approved sources...'})}\n\n"
        else:
            yield f"data: {json.dumps({'status': 'researching', 'message': f'Agent searching for real sources on {topic}...'})}\n\n"
        full_text = ""
        try:
            if AI_PROVIDER == "gemini":
                model = genai.GenerativeModel("gemini-2.0-flash-lite")
                response = model.generate_content(prompt, stream=True)
                for chunk in response:
                    if chunk.text:
                        full_text += chunk.text
                        yield f"data: {json.dumps({'text': chunk.text})}\n\n"
            else:
                response = openai_client.chat.completions.create(
                    model="gpt-4o-mini",
                    messages=[{"role": "user", "content": prompt}],
                    stream=True,
                )
                for chunk in response:
                    delta = chunk.choices[0].delta.content or ""
                    if delta:
                        full_text += delta
                        yield f"data: {json.dumps({'text': delta})}\n\n"
        except Exception as e:
            print(f"Stream error: {e}")
            yield "data: [DONE]\n\n"
            return
        print(f"Stream complete, full_text length: {len(full_text)}")

        def parse_curriculum(text):
            clean = text.replace("```json\n", "").replace("```\n", "").replace("```", "").strip()
            first = clean.index("{")
            last = clean.rindex("}")
            return json.loads(clean[first:last + 1])

        def validate_structure(parsed, expected_count):
            """Check complexity_level progression and module count."""
            modules = parsed.get("modules", [])
            if len(modules) != expected_count:
                return False, f"Expected {expected_count} modules, got {len(modules)}"
            levels = [m.get("complexity_level", 0) for m in modules]
            if levels[0] != 1:
                return False, f"First module complexity should be 1, got {levels[0]}"
            if levels[-1] != 5:
                return False, f"Last module complexity should be 5, got {levels[-1]}"
            for i in range(1, len(levels)):
                if levels[i] < levels[i-1]:
                    return False, f"Complexity decreased at module {i+1}"
            return True, "ok"

        # Save to DB BEFORE sending [DONE] — client disconnects on [DONE]
        parsed = None
        try:
            parsed = parse_curriculum(full_text)
            valid, reason = validate_structure(parsed, module_count)
            if not valid:
                print(f"Validation failed: {reason} — retrying once")
                yield f"data: {json.dumps({'status': 'fixing', 'message': f'Fixing structure: {reason}...'})} \n\n"
                fix_prompt = prompt + f"\n\nIMPORTANT: Your previous response had a structural error: {reason}. Fix it and return valid JSON only."
                if AI_PROVIDER == "gemini":
                    model = genai.GenerativeModel("gemini-2.0-flash-lite")
                    retry_response = model.generate_content(fix_prompt)
                    retry_text = retry_response.text
                else:
                    retry_response = openai_client.chat.completions.create(
                        model="gpt-4o-mini",
                        messages=[{"role": "user", "content": fix_prompt}],
                    )
                    retry_text = retry_response.choices[0].message.content
                try:
                    parsed = parse_curriculum(retry_text)
                    retry_valid, retry_reason = validate_structure(parsed, module_count)
                    yield f"data: {json.dumps({'reset': True})}\n\n"
                    yield f"data: {json.dumps({'text': retry_text})}\n\n"
                    if retry_valid:
                        print("Retry succeeded")
                    else:
                        actual_count = len(parsed.get("modules", []))
                        print(f"Retry also failed validation: {retry_reason}")
                        yield f"data: {json.dumps({'type': 'warning', 'message': f'Generated {actual_count} modules instead of {module_count} — GPT was being lazy, try regenerating'})}\n\n"
                except Exception as e:
                    print(f"Retry parse failed: {e}")
            for m in (parsed.get("modules") or []):
                m["learning_objectives"] = [
                    o[0].upper() + o[1:] if o else o for o in (m.get("learning_objectives") or [])
                ]
            save_curriculum(topic, level, audience, course_code, course_type, module_count, parsed, design_approach)
            print(f"Saved curriculum: {topic}")
        except Exception as e:
            print(f"Failed to parse/save curriculum: {e}")
        yield "data: [DONE]\n\n"

    return Response(
        stream_with_context(event_stream()),
        mimetype="text/event-stream",
        headers={
            "Cache-Control": "no-cache",
            "X-Accel-Buffering": "no",
        },
    )


@app.route("/api/curriculum/save", methods=["POST"])
def save_curriculum_endpoint():
    """Save a fully expanded curriculum (from two-phase generation) to history."""
    data = request.get_json()
    topic = data.get("topic", "")
    level = data.get("level", "")
    audience = data.get("audience", "")
    course_code = data.get("course_code", "")
    course_type = data.get("course_type", "mixed")
    module_count = data.get("module_count", 0)
    design_approach = data.get("design_approach", "ADDIE")
    modules = data.get("modules", [])
    sources = data.get("sources", [])
    course_narrative = data.get("course_narrative", "")
    parsed = {"modules": modules, "sources": sources, "course_narrative": course_narrative}
    try:
        save_curriculum(topic, level, audience, course_code, course_type, module_count, parsed, design_approach)
        return jsonify({"status": "saved"})
    except Exception as e:
        print(f"Save endpoint error: {e}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/curriculum/skeleton", methods=["POST"])
def generate_skeleton():
    """Phase 1: Generate only module titles + learning_objectives (no readings/assignments)."""
    data = request.get_json()
    topic = data.get("topic", "")
    level = data.get("level", "")
    audience = data.get("audience", "")
    accreditation_context = data.get("accreditation_context", "")
    course_code = data.get("course_code", "")
    course_type = data.get("course_type", "mixed")
    module_count_raw = data.get("module_count", "6")
    design_approach = data.get("design_approach", "addie").lower()
    if design_approach not in ("addie", "sam"):
        design_approach = "addie"

    if not all([topic, level, audience]):
        return {"error": "Missing required fields"}, 400

    try:
        module_count = max(3, min(12, int(module_count_raw)))
    except (ValueError, TypeError):
        module_count = 6

    blooms_constraint = get_blooms_constraint(level)

    prompt = f"""You are an expert curriculum designer. Generate ONLY the module skeleton for the following course.

Topic: {topic}
Course Code: {course_code or "Not specified"}
Level: {level}
Target Audience: {audience}
Accreditation Context: {accreditation_context or "None"}
Course Type: {course_type}
Number of Modules: {module_count}
Design Approach: {design_approach}

Bloom's Verb Constraint: {blooms_constraint}
Difficulty Progression: complexity_level must start at 1 and reach 5 by the final module, increasing evenly.

Generate the course skeleton. Include a course_narrative (2-3 sentences explaining the central question or theme of this course and why these modules belong together — the "story" of the whole course). For each module provide: module_number, title, complexity_level, learning_objectives (list of 2-3 objectives using the permitted Bloom's verbs). Nothing else — no readings, no assignments, no narrative_preview.

Return ONLY valid JSON (no markdown, no explanation):
{{
  "course_narrative": "A 2-3 sentence explanation of the course's central theme and why these modules belong together.",
  "modules": [
    {{
      "module_number": 1,
      "title": "Module title",
      "complexity_level": 1,
      "learning_objectives": ["objective using permitted Bloom's verbs", "objective 2"]
    }}
  ]
}}

CRITICAL: Generate exactly {module_count} modules. complexity_level must start at 1 and end at 5."""

    def event_stream():
        yield f"data: {json.dumps({'status': 'generating', 'message': f'Generating {module_count}-module skeleton for {topic}...'})}\n\n"
        full_text = ""
        try:
            if AI_PROVIDER == "gemini":
                model = genai.GenerativeModel("gemini-2.0-flash-lite")
                response = model.generate_content(prompt, stream=True)
                for chunk in response:
                    if chunk.text:
                        full_text += chunk.text
                        yield f"data: {json.dumps({'text': chunk.text})}\n\n"
            else:
                response = openai_client.chat.completions.create(
                    model="gpt-4o-mini",
                    messages=[{"role": "user", "content": prompt}],
                    stream=True,
                )
                for chunk in response:
                    delta = chunk.choices[0].delta.content or ""
                    if delta:
                        full_text += delta
                        yield f"data: {json.dumps({'text': delta})}\n\n"
        except Exception as e:
            print(f"Skeleton stream error: {e}")
            yield "data: [DONE]\n\n"
            return

        # Validate skeleton structure
        def parse_skeleton(text):
            clean = text.replace("```json\n", "").replace("```\n", "").replace("```", "").strip()
            first = clean.index("{")
            last = clean.rindex("}")
            return json.loads(clean[first:last + 1])

        try:
            parsed = parse_skeleton(full_text)
            modules = parsed.get("modules", [])
            # Ensure module_number field exists
            for i, m in enumerate(modules):
                if "module_number" not in m:
                    m["module_number"] = i + 1
            # Validate count
            if len(modules) != module_count:
                print(f"Skeleton count mismatch: expected {module_count}, got {len(modules)}")
        except Exception as e:
            print(f"Failed to parse skeleton: {e}")

        yield "data: [DONE]\n\n"

    return Response(
        stream_with_context(event_stream()),
        mimetype="text/event-stream",
        headers={
            "Cache-Control": "no-cache",
            "X-Accel-Buffering": "no",
        },
    )


@app.route("/api/curriculum/expand", methods=["POST"])
def expand_module():
    """Phase 2: Expand a single skeleton module with readings, assignments, narrative, etc."""
    data = request.get_json()
    skeleton = data.get("skeleton", [])       # full modules array from skeleton phase
    module_index = data.get("module_index", 0)
    topic = data.get("topic", "")
    level = data.get("level", "")
    audience = data.get("audience", "")
    course_type = data.get("course_type", "mixed")
    design_approach = data.get("design_approach", "addie").lower()
    course_code = data.get("course_code", "")
    accreditation_context = data.get("accreditation_context", "")
    approved_sources_raw = data.get("approved_sources", [])

    try:
        session_duration = max(1, int(data.get("session_duration", 90)))
    except (ValueError, TypeError):
        session_duration = 90

    if not all([topic, level, audience]) or not skeleton or module_index >= len(skeleton):
        return {"error": "Missing required fields or invalid module_index"}, 400

    module = skeleton[module_index]
    module_title = module.get("title", f"Module {module_index + 1}")
    module_number = module.get("module_number", module_index + 1)
    complexity_level = module.get("complexity_level", 1)
    learning_objectives = module.get("learning_objectives", [])
    total_modules = len(skeleton)

    blooms_constraint = get_blooms_constraint(level)
    session_constraint = get_session_constraints(session_duration)
    assessment_format = ASSESSMENT_FORMATS.get(course_type, ASSESSMENT_FORMATS["mixed"])

    # Build sources context from approved sources
    sources_context = ""
    required_sources = []
    optional_sources = []
    if approved_sources_raw and isinstance(approved_sources_raw, list):
        real_sources = [s for s in approved_sources_raw if s.get("url")]
        if real_sources:
            sources_context = "\n\nApproved sources from instructor — use these URLs for readings where relevant:\n"
            for s in real_sources:
                priority = s.get("priority", "optional")
                tag = f" [REQUIRED]" if priority == "required" else " [OPTIONAL]"
                sources_context += f"- [{s.get('type', 'other')}]{tag} {s.get('title', '')} | {s.get('url', '')}\n"
                if priority == "required":
                    required_sources.append(s)
                else:
                    optional_sources.append(s)
            sources_context += "Prioritize required sources. Do NOT invent URLs.\n"

    if design_approach == "sam":
        design_approach_label = "SAM (Successive Approximation Model)"
        sam_field = '"rapid_prototype_cycle": "Description of the Rapid Prototype → Evaluate → Revise cycle for this module.",'
    else:
        design_approach_label = "ADDIE (Analysis → Design → Development → Implementation → Evaluation)"
        sam_field = ""

    # Resource type priority by course type
    resource_priority_map = {
        "project": "PRIORITY: recommended_readings should lean toward academic + news sources (real-world cases and current research). Include at least 1 news/industry source per module where relevant.",
        "essay": "PRIORITY: recommended_readings should lean toward academic sources. Include video (TED/lecture) and news where they support the argument. Minimum 1 academic per module.",
        "debate": "PRIORITY: recommended_readings should include news/current events AND academic sources to support multiple perspectives. At least 1 news source per module.",
        "lab": "PRIORITY: recommended_readings should lean heavily toward video resources (tutorials, demonstrations, walkthroughs). At least 1 video per module where possible.",
        "mixed": "PRIORITY: recommended_readings should include a balanced mix of academic, news, and video sources across modules.",
    }
    resource_priority = resource_priority_map.get(course_type, resource_priority_map["mixed"])

    objectives_str = "\n".join(f"  - {obj}" for obj in learning_objectives)

    prompt = f"""You are an expert curriculum designer. Expand the following module skeleton into a full module with all required fields.

Course Context:
- Topic: {topic}
- Course Code: {course_code or "Not specified"}
- Level: {level}
- Target Audience: {audience}
- Accreditation: {accreditation_context or "None"}
- Course Type: {course_type}
- Design Approach: {design_approach_label}
- Total Modules in Course: {total_modules}

Module to Expand:
- Module Number: {module_number} of {total_modules}
- Title: {module_title}
- Complexity Level: {complexity_level}/5
- Learning Objectives:
{objectives_str}

Constraints:
- Bloom's Verb Constraint: {blooms_constraint}
- Session Duration: {session_constraint}
- Assessment Format: {assessment_format}
- Max 2 recommended readings. Each must have a clear rationale tied to this module's learning objectives.
- {resource_priority}
- Assignment task_description: MUST be specific and actionable (e.g. "Write a 500-word reflection comparing two case studies..."), NOT generic.
- Assignment rubric_highlights: MUST contain exactly 3-4 concrete criteria.
- Not every module requires an assignment. Only include one if it meaningfully fits this module.
- narrative_preview: A compelling 2-3 sentence narrative hook using metaphor, scenario, or challenge framing.{sources_context}

Return ONLY valid JSON for this single module (no markdown, no explanation):
{{
  "module_number": {module_number},
  "title": {json.dumps(module_title)},
  "complexity_level": {complexity_level},
  "learning_objectives": {json.dumps(learning_objectives)},
  {sam_field}
  "narrative_preview": "A compelling 2-3 sentence narrative hook.",
  "recommended_readings": [
    {{
      "title": "Full title of reading",
      "url": "https://real-url.com",
      "type": "academic | video | news",
      "estimated_time": "15 min read",
      "reading_type": "required | optional",
      "key_points": ["key point 1", "key point 2"],
      "rationale": "Why this reading is essential for this module."
    }}
  ],
  "assignments": [
    {{
      "type": "project | essay | quiz | discussion | presentation | lab | reflection",
      "title": "Short assignment title",
      "task_description": "2-3 sentence specific description of exactly what students must do.",
      "deliverable": "What they hand in",
      "estimated_time": "Realistic time given session duration",
      "covers_objectives": "Which specific learning objectives this addresses",
      "rubric_highlights": [
        "Criterion 1 — description of excellent work",
        "Criterion 2 — description of excellent work",
        "Criterion 3 — description of excellent work"
      ]
    }}
  ]
}}"""

    def event_stream():
        yield f"data: {json.dumps({'status': 'expanding', 'message': f'Expanding module {module_number}: {module_title}...'})}\n\n"
        full_text = ""
        try:
            if AI_PROVIDER == "gemini":
                model = genai.GenerativeModel("gemini-2.0-flash-lite")
                response = model.generate_content(prompt, stream=True)
                for chunk in response:
                    if chunk.text:
                        full_text += chunk.text
                        yield f"data: {json.dumps({'text': chunk.text})}\n\n"
            else:
                response = openai_client.chat.completions.create(
                    model="gpt-4o-mini",
                    messages=[{"role": "user", "content": prompt}],
                    stream=True,
                )
                for chunk in response:
                    delta = chunk.choices[0].delta.content or ""
                    if delta:
                        full_text += delta
                        yield f"data: {json.dumps({'text': delta})}\n\n"
        except Exception as e:
            print(f"Expand module stream error (module {module_index}): {e}")
            yield "data: [DONE]\n\n"
            return

        # Validate expanded module
        try:
            clean = full_text.replace("```json\n", "").replace("```\n", "").replace("```", "").strip()
            first = clean.index("{")
            last = clean.rindex("}")
            parsed = json.loads(clean[first:last + 1])
            # Ensure required array fields
            if not isinstance(parsed.get("recommended_readings"), list):
                parsed["recommended_readings"] = []
            if not isinstance(parsed.get("assignments"), list):
                parsed["assignments"] = []
            if not isinstance(parsed.get("learning_objectives"), list):
                parsed["learning_objectives"] = learning_objectives
            parsed["learning_objectives"] = [
                o[0].upper() + o[1:] if o else o for o in parsed["learning_objectives"]
            ]
            print(f"Expanded module {module_number}: {module_title}")
        except Exception as e:
            print(f"Failed to parse expanded module {module_index}: {e}")

        yield "data: [DONE]\n\n"

    return Response(
        stream_with_context(event_stream()),
        mimetype="text/event-stream",
        headers={
            "Cache-Control": "no-cache",
            "X-Accel-Buffering": "no",
        },
    )


@app.route("/api/history", methods=["GET"])
def get_history():
    conn = get_db()
    if not conn:
        return {"history": [], "error": "DB unavailable"}, 200
    try:
        cur = conn.cursor()
        cur.execute(
            "SELECT id, created_at, topic, level, course_code, course_type, module_count, is_favorite "
            "FROM curricula ORDER BY is_favorite DESC, created_at DESC LIMIT 50"
        )
        rows = cur.fetchall()
        cur.close()
        conn.close()
        return {"history": [
            {
                "id": r[0],
                "created_at": r[1].isoformat(),
                "topic": r[2],
                "level": r[3],
                "course_code": r[4] or "",
                "course_type": r[5] or "mixed",
                "module_count": r[6],
                "is_favorite": r[7] or False,
            }
            for r in rows
        ]}
    except Exception as e:
        return {"history": [], "error": str(e)}, 200


@app.route("/api/history/<int:curriculum_id>", methods=["GET"])
def get_curriculum_by_id(curriculum_id):
    conn = get_db()
    if not conn:
        return {"error": "DB unavailable"}, 503
    try:
        cur = conn.cursor()
        cur.execute(
            "SELECT topic, level, audience, course_code, course_type, modules, sources "
            "FROM curricula WHERE id = %s",
            (curriculum_id,)
        )
        row = cur.fetchone()
        cur.close()
        conn.close()
        if not row:
            return {"error": "Not found"}, 404
        # modules/sources are stored via json.dumps so psycopg2 may return
        # them as a raw string (if JSONB oid decoding is not registered).
        # Normalise to Python objects defensively.
        def _decode(val):
            if isinstance(val, str):
                try:
                    return json.loads(val)
                except Exception:
                    return []
            return val if val is not None else []
        modules = _decode(row[5])
        sources = _decode(row[6])
        # Ensure every module has the expected array fields so the frontend
        # never calls .map() on null/undefined.
        for m in (modules if isinstance(modules, list) else []):
            if not isinstance(m.get("learning_objectives"), list):
                m["learning_objectives"] = []
            if not isinstance(m.get("recommended_readings"), list):
                m["recommended_readings"] = []
            if not isinstance(m.get("assignments"), list):
                m["assignments"] = []
        return {
            "topic": row[0],
            "level": row[1],
            "audience": row[2],
            "course_code": row[3] or "",
            "course_type": row[4] or "mixed",
            "modules": modules,
            "sources": sources,
        }
    except Exception as e:
        return {"error": str(e)}, 500


@app.route("/api/history/<int:curriculum_id>", methods=["DELETE"])
def delete_curriculum(curriculum_id):
    conn = get_db()
    if not conn:
        return {"error": "DB unavailable"}, 503
    try:
        cur = conn.cursor()
        cur.execute("DELETE FROM curricula WHERE id = %s", (curriculum_id,))
        conn.commit()
        cur.close()
        conn.close()
        return {"status": "deleted"}
    except Exception as e:
        return {"error": str(e)}, 500


@app.route("/api/history/<int:curriculum_id>/favorite", methods=["POST"])
def toggle_favorite(curriculum_id):
    conn = get_db()
    if not conn:
        return {"error": "DB unavailable"}, 503
    try:
        cur = conn.cursor()
        cur.execute(
            "UPDATE curricula SET is_favorite = NOT is_favorite WHERE id = %s RETURNING is_favorite",
            (curriculum_id,)
        )
        row = cur.fetchone()
        conn.commit()
        cur.close()
        conn.close()
        return {"is_favorite": row[0] if row else False}
    except Exception as e:
        return {"error": str(e)}, 500


@app.route("/api/xapi/statement", methods=["POST"])
def receive_xapi():
    statement = request.get_json()
    actor = statement.get("actor", {}).get("name", "unknown")
    verb = statement.get("verb", {}).get("display", {}).get("en-US", "unknown")
    obj = statement.get("object", {}).get("definition", {}).get("name", {}).get("en-US", "unknown")
    print(f"xAPI: {actor} {verb} {obj}")
    return {"status": "received"}, 200


def _get_lightrag_instance(storage_dir: str = None):
    """Return a cached LightRAG instance (not yet initialized — init happens inside async context)."""
    from lightrag import LightRAG
    from lightrag.llm.openai import gpt_4o_mini_complete, openai_embed
    from lightrag.utils import EmbeddingFunc

    if storage_dir is None:
        backend_dir = os.path.dirname(os.path.abspath(__file__))
        storage_dir = os.path.normpath(os.path.join(backend_dir, "..", "data", "lightrag_storage"))

    if storage_dir in _rag_instances:
        return _rag_instances[storage_dir]

    rag = LightRAG(
        working_dir=storage_dir,
        llm_model_func=gpt_4o_mini_complete,
        embedding_func=EmbeddingFunc(
            embedding_dim=1536,
            max_token_size=8192,
            func=lambda texts: openai_embed(texts, model="text-embedding-3-small"),
        ),
    )
    _rag_instances[storage_dir] = rag
    return rag


def _get_graphml_path(subject: str = "all") -> str | None:
    """Return the path to the graphml file for a specific (non-all) subject.

    Returns None for unknown/unrecognised subject keys so the caller can treat
    them as 'not_ready' rather than silently falling back to the Business Law
    storage directory.
    """
    backend_dir = os.path.dirname(os.path.abspath(__file__))
    if subject == "call":
        storage_dir = "lightrag_storage_call"
    elif subject == "business-law":
        storage_dir = "lightrag_storage"
    else:
        # Dynamic subject — storage dir follows the convention used at ingest time.
        storage_dir = f"lightrag_storage_{subject}"
    return os.path.normpath(
        os.path.join(backend_dir, "..", "data", storage_dir, "graph_chunk_entity_relation.graphml")
    )


def _get_all_graphml_paths() -> list:
    """Return paths for both graph files used in the 'all' merged view."""
    backend_dir = os.path.dirname(os.path.abspath(__file__))
    data_dir = os.path.normpath(os.path.join(backend_dir, "..", "data"))
    return [
        os.path.join(data_dir, "lightrag_storage", "graph_chunk_entity_relation.graphml"),
        os.path.join(data_dir, "lightrag_storage_call", "graph_chunk_entity_relation.graphml"),
    ]


def _parse_graph_from_file(graphml_path: str):
    """Read a graphml file and return (nodes_dict, edges_set_data).

    nodes_dict  → {node_id: attrs_dict}
    edges_list  → list of (source, target, attrs_dict)
    """
    import networkx as nx
    G = nx.read_graphml(graphml_path)
    return G


def _build_graph_response(graphs):
    """Merge one or more networkx graphs and return filtered {nodes, edges} dicts.

    Deduplication rules:
    - Nodes: keyed by node ID. If the same ID appears in multiple graphs, keep the
      version with more attributes (len(attrs)).
    - Edges: keyed by (source, target, relation). First occurrence wins.
    """
    PERSON_TYPES = {"person", "PERSON"}
    PERSON_DESC_PHRASES = ("a person", "a student", "a fictional")

    merged_nodes = {}  # node_id → attrs dict
    merged_edges = {}  # (source, target, relation) → attrs dict

    for G in graphs:
        for node_id, attrs in G.nodes(data=True):
            nid = str(node_id)
            if nid not in merged_nodes or len(attrs) > len(merged_nodes[nid]):
                merged_nodes[nid] = dict(attrs)

        for source, target, attrs in G.edges(data=True):
            relation = attrs.get("relation", attrs.get("label", ""))
            key = (str(source), str(target), relation)
            if key not in merged_edges:
                merged_edges[key] = dict(attrs)

    # Date-noise patterns to suppress
    _DATE_FULL_RE = re.compile(
        r'^(January|February|March|April|May|June|July|August|September|October|November|December)'
        r'\s+\d{1,2},?\s+\d{4}$',
        re.IGNORECASE,
    )
    _DATE_YEAR_RE = re.compile(r'^\d{4}$')

    # Filter PERSON and date nodes
    filtered_node_ids = set()
    nodes = []
    for node_id, attrs in merged_nodes.items():
        entity_type = attrs.get("entity_type", "")
        raw_desc = attrs.get("description", "")
        if raw_desc and "<SEP>" in raw_desc:
            raw_desc = raw_desc.split("<SEP>")[0].strip()

        if entity_type in PERSON_TYPES:
            filtered_node_ids.add(node_id)
            continue
        if raw_desc and any(phrase in raw_desc.lower() for phrase in PERSON_DESC_PHRASES):
            filtered_node_ids.add(node_id)
            continue

        # Filter date-only nodes (e.g. "January 24, 2025" or "2025")
        node_label = attrs.get("label", node_id)
        if _DATE_FULL_RE.match(str(node_label)) or _DATE_YEAR_RE.match(str(node_label)):
            filtered_node_ids.add(node_id)
            continue

        nodes.append({
            "id": node_id,
            "label": node_label,
            "entity_type": entity_type,
            "description": raw_desc,
        })

    edges = []
    for (source, target, relation), attrs in merged_edges.items():
        if source in filtered_node_ids or target in filtered_node_ids:
            continue
        edges.append({
            "source": source,
            "target": target,
            "label": attrs.get("label", relation),
        })

    return nodes, edges


@app.route("/api/graph", methods=["GET"])
def get_graph():
    """Return knowledge graph nodes and edges from the LightRAG graphml file(s).

    subject=all (default) → merge business-law + CALL graphs
    subject=business-law  → business-law graph only
    subject=call          → CALL graph only
    """
    import networkx as nx
    subject = request.args.get("subject", "all")

    try:
        if subject == "all":
            paths = _get_all_graphml_paths()
            existing_paths = [p for p in paths if os.path.exists(p)]
            if not existing_paths:
                return {"nodes": [], "edges": [], "status": "not_ready"}
            graphs = [nx.read_graphml(p) for p in existing_paths]
        else:
            graphml_path = _get_graphml_path(subject)
            if graphml_path is None or not os.path.exists(graphml_path):
                return {"nodes": [], "edges": [], "status": "not_ready"}
            graphs = [nx.read_graphml(graphml_path)]

        nodes, edges = _build_graph_response(graphs)
        return {"nodes": nodes, "edges": edges, "status": "ready"}

    except Exception as e:
        print(f"Graph read error: {e}")
        return {"nodes": [], "edges": [], "status": "error", "error": str(e)}, 500


@app.route("/api/graph/query", methods=["POST"])
def query_graph():
    """Query the LightRAG knowledge graph and return an answer.

    Accepts optional 'subject' field (default: 'business-law') to select which
    knowledge graph to query. Uses:
    - Layer A: module-level _rag_instances dict to avoid re-initializing LightRAG
    - Layer B: Redis cache (TTL 24h) keyed on subject + normalized question
    """
    data = request.get_json()
    question = data.get("question", "").strip()
    mode = data.get("mode", "hybrid")
    subject = data.get("subject", "business-law")

    if not question:
        return {"error": "Missing 'question' field."}, 400

    graphml_path = _get_graphml_path(subject)
    if graphml_path is None or not os.path.exists(graphml_path):
        return {"answer": "Knowledge graph not initialized yet.", "subject": subject, "matched_node_id": None}

    def clean_answer(text: str) -> str:
        """Strip markdown formatting and truncate to 3 sentences."""
        text = re.sub(r'#{1,6}\s*', '', text)
        text = re.sub(r'\*\*(.*?)\*\*', r'\1', text)
        text = re.sub(r'\*(.*?)\*', r'\1', text)
        text = re.sub(r'\[\d+\]', '', text)
        text = re.sub(r'(?m)^\s*[-*•]\s+', '', text)
        text = re.sub(r'\n{2,}', ' ', text)
        text = re.sub(r'[ \t]+', ' ', text)
        text = text.strip()
        sentences = re.split(r'(?<=[.!?])\s+', text)
        sentences = [s.strip() for s in sentences if s.strip() and len(s.strip()) > 10]
        result = ' '.join(sentences[:3])
        if result and result[-1] not in '.!?':
            result += '.'
        return result

    # --- Layer B: Redis cache check ---
    normalized_q = question.lower().strip()
    cache_key = f"graph_query:{subject}:{normalized_q}"
    if _redis_client is not None:
        try:
            cached = _redis_client.get(cache_key)
            if cached:
                print(f"Redis cache hit: {cache_key}")
                # Still resolve matched_node_id from the graph even on cache hit
                cached_node_id = None
                try:
                    import networkx as nx
                    _G = nx.read_graphml(graphml_path)
                    _q = question.lower().strip()
                    for _nid, _attrs in _G.nodes(data=True):
                        _lbl = _attrs.get("label", str(_nid)).lower()
                        if _lbl == _q or _lbl.startswith(_q) or _q in _lbl:
                            cached_node_id = str(_nid)
                            break
                except Exception:
                    pass
                return {"answer": cached, "subject": subject, "matched_node_id": cached_node_id, "cached": True}
        except Exception as redis_err:
            print(f"Redis get error (skipping cache): {redis_err}")

    try:
        import asyncio
        try:
            from lightrag import QueryParam
        except ImportError:
            return {"answer": "Query engine not available in this environment. Run the backend locally with lightrag installed.", "subject": subject, "matched_node_id": None}

        # --- Layer A: use cached LightRAG instance (initialize_storages only called once) ---
        backend_dir = os.path.dirname(os.path.abspath(__file__))
        if subject == "call":
            storage_dir = os.path.normpath(os.path.join(backend_dir, "..", "data", "lightrag_storage_call"))
        else:
            storage_dir = os.path.normpath(os.path.join(backend_dir, "..", "data", "lightrag_storage"))

        # Expand abbreviation queries (e.g. "DGBLL" → "Digital Game-Based Language Learning (DGBLL)")
        def _expand_abbreviation(q: str, graphml_path: str) -> str:
            q_stripped = q.strip()
            try:
                import networkx as nx
                G = nx.read_graphml(graphml_path)
                node_labels = [attrs.get("label", str(node_id)) for node_id, attrs in G.nodes(data=True)]

                # Check if query already matches a node exactly — no expansion needed
                if q_stripped in node_labels:
                    return q_stripped

                q_lower = q_stripped.lower()

                # 1. Look for nodes that START WITH the query (case-insensitive)
                starts_with = [n for n in node_labels if n.lower().startswith(q_lower)]
                if starts_with:
                    best = min(starts_with, key=len)
                    return f"What is {best}? {q_stripped}"

                # 2. Look for nodes that CONTAIN the query as a word
                contains = [n for n in node_labels if q_lower in n.lower()]
                if contains:
                    best = min(contains, key=len)
                    return f"What is {best}? {q_stripped}"

                # 3. Original parenthesis match as fallback
                for n in node_labels:
                    if f"({q_stripped})" in n or f"({q_stripped.upper()})" in n:
                        return f"What is {n}? {q_stripped}"

            except Exception:
                pass
            return q_stripped

        expanded_question = _expand_abbreviation(question, graphml_path)

        async def _run_query():
            rag = _get_lightrag_instance(storage_dir)
            if storage_dir not in _initialized_instances:
                await rag.initialize_storages()
                _initialized_instances.add(storage_dir)
            return await rag.aquery(expanded_question, param=QueryParam(mode=mode))

        raw_answer = _run_async(_run_query())

        answer = clean_answer(raw_answer)

        # --- Find best matching node for highlight (runs for ALL queries) ---
        NO_INFO_PHRASES = [
            "not have enough information",
            "don't have enough information",
            "cannot answer",
            "no information",
        ]
        matched_node_id = None
        answer_lower = answer.lower()
        is_no_info = any(phrase in answer_lower for phrase in NO_INFO_PHRASES)

        try:
            import networkx as nx
            G = nx.read_graphml(graphml_path)
            q_lower = question.lower()

            best_node_id = None
            best_node_label = None
            best_node_desc = None

            # Priority 1: exact label match on original question
            for node_id, attrs in G.nodes(data=True):
                label = attrs.get("label", str(node_id))
                if label.lower() == q_lower:
                    best_node_id = node_id
                    best_node_label = label
                    best_node_desc = attrs.get("description", "")
                    break

            # Priority 2: starts-with match (shortest wins)
            if best_node_id is None:
                candidates = [
                    (node_id, attrs.get("label", str(node_id)), attrs.get("description", ""))
                    for node_id, attrs in G.nodes(data=True)
                    if attrs.get("label", str(node_id)).lower().startswith(q_lower)
                ]
                if candidates:
                    best_node_id, best_node_label, best_node_desc = min(candidates, key=lambda x: len(x[1]))

            # Priority 3: contains match (shortest wins)
            if best_node_id is None:
                candidates = [
                    (node_id, attrs.get("label", str(node_id)), attrs.get("description", ""))
                    for node_id, attrs in G.nodes(data=True)
                    if q_lower in attrs.get("label", str(node_id)).lower()
                ]
                if candidates:
                    best_node_id, best_node_label, best_node_desc = min(candidates, key=lambda x: len(x[1]))

            if best_node_id is not None:
                matched_node_id = str(best_node_id)
                # Only override answer if LightRAG had no info
                if is_no_info and best_node_desc:
                    answer = f"Based on the knowledge graph: {best_node_desc}"
        except Exception as fallback_err:
            print(f"Node matching error: {fallback_err}")

        # --- Layer B: store result in Redis ---
        if _redis_client is not None:
            try:
                _redis_client.set(cache_key, answer)  # permanent — graph doesn't change
            except Exception as redis_err:
                print(f"Redis set error (skipping cache store): {redis_err}")

        return {"answer": answer, "subject": subject, "matched_node_id": matched_node_id}
    except Exception as e:
        print(f"Graph query error: {e}")
        return {"answer": f"Query failed: {str(e)}"}, 500


def seed_mock_xapi():
    """Seed mock xAPI statements if table is empty."""
    conn = get_db()
    if not conn:
        return
    cur = conn.cursor()
    cur.execute("SELECT COUNT(*) FROM xapi_statements")
    count = cur.fetchone()[0]
    if count > 0:
        conn.close()
        return

    from datetime import datetime, timedelta
    import random

    statements = [
        # Alice — high performer, completed most things
        ("alice@test.com", "Alice Chen", "completed", "module/1", "Introduction to the Course"),
        ("alice@test.com", "Alice Chen", "passed", "module/1/reading/0", "Week 1 Reading"),
        ("alice@test.com", "Alice Chen", "completed", "module/2", "Core Concepts"),
        ("alice@test.com", "Alice Chen", "attempted", "module/3", "Applied Theory"),
        ("alice@test.com", "Alice Chen", "passed", "module/3/quiz", "Module 3 Quiz"),
        # Bob — struggling on module 3
        ("bob@test.com", "Bob Kim", "completed", "module/1", "Introduction to the Course"),
        ("bob@test.com", "Bob Kim", "experienced", "module/2/reading/0", "Core Reading"),
        ("bob@test.com", "Bob Kim", "struggled", "module/3", "Applied Theory"),
        ("bob@test.com", "Bob Kim", "struggled", "module/3/concept/theory", "Theoretical Framework"),
        ("bob@test.com", "Bob Kim", "attempted", "module/3/quiz", "Module 3 Quiz"),
        # Carol — dropped off after module 2
        ("carol@test.com", "Carol Singh", "experienced", "module/1", "Introduction to the Course"),
        ("carol@test.com", "Carol Singh", "completed", "module/1", "Introduction to the Course"),
        ("carol@test.com", "Carol Singh", "experienced", "module/2", "Core Concepts"),
        ("carol@test.com", "Carol Singh", "struggled", "module/2/concept/advanced", "Advanced Core Concept"),
        # David — steady progress
        ("david@test.com", "David Park", "completed", "module/1", "Introduction to the Course"),
        ("david@test.com", "David Park", "completed", "module/2", "Core Concepts"),
        ("david@test.com", "David Park", "passed", "module/2/quiz", "Module 2 Quiz"),
        ("david@test.com", "David Park", "attempted", "module/3", "Applied Theory"),
        ("david@test.com", "David Park", "experienced", "module/4", "Case Studies"),
        ("david@test.com", "David Park", "struggled", "module/4/concept/integration", "Integration Concept"),
    ]

    base_time = datetime.now() - timedelta(days=7)
    for i, (email, name, verb, obj_id, obj_name) in enumerate(statements):
        ts = base_time + timedelta(hours=i * 3 + random.randint(0, 2))
        cur.execute(
            "INSERT INTO xapi_statements (actor_email, actor_name, verb, object_id, object_name, timestamp, curriculum_topic) VALUES (%s, %s, %s, %s, %s, %s, %s)",
            (email, name, verb, obj_id, obj_name, ts, "Mock Course")
        )
    conn.commit()
    conn.close()
    print("Mock xAPI statements seeded.")


@app.route("/xapi/statements", methods=["POST"])
def receive_xapi_statement():
    """Receive a single xAPI statement and store it."""
    data = request.get_json()
    actor = data.get("actor", {})
    verb = data.get("verb", {})
    obj = data.get("object", {})

    email = actor.get("mbox", "").replace("mailto:", "")
    name = actor.get("name", email)
    verb_id = verb.get("id", "").split("/")[-1]
    obj_id = obj.get("id", "")
    obj_name = obj.get("definition", {}).get("name", {}).get("en-US", obj_id)
    curriculum_topic = data.get("context", {}).get("extensions", {}).get("curriculum_topic", "")

    conn = get_db()
    if not conn:
        return jsonify({"error": "DB unavailable"}), 503
    cur = conn.cursor()
    cur.execute(
        "INSERT INTO xapi_statements (actor_email, actor_name, verb, object_id, object_name, curriculum_topic) VALUES (%s, %s, %s, %s, %s, %s)",
        (email, name, verb_id, obj_id, obj_name, curriculum_topic)
    )
    conn.commit()
    conn.close()

    # Update Redis learner state
    if _redis_client:
        key = f"learner:{email}"
        state = _redis_client.get(key)
        import json as _json
        state_data = _json.loads(state) if state else {"viewed": [], "struggling": [], "mastered": []}
        if verb_id == "experienced" and obj_id not in state_data["viewed"]:
            state_data["viewed"].append(obj_id)
        elif verb_id == "struggled" and obj_id not in state_data["struggling"]:
            state_data["struggling"].append(obj_id)
        elif verb_id in ("completed", "passed") and obj_id not in state_data["mastered"]:
            state_data["mastered"].append(obj_id)
        _redis_client.set(key, _json.dumps(state_data))

    return jsonify({"status": "stored"})


@app.route("/xapi/statements", methods=["GET"])
def get_xapi_statements():
    """Return recent xAPI statements (last 50)."""
    conn = get_db()
    if not conn:
        return jsonify([])
    cur = conn.cursor()
    cur.execute(
        "SELECT actor_name, actor_email, verb, object_name, object_id, timestamp, curriculum_topic FROM xapi_statements ORDER BY timestamp DESC LIMIT 50"
    )
    rows = cur.fetchall()
    conn.close()
    statements = [
        {"actor_name": r[0], "actor_email": r[1], "verb": r[2], "object_name": r[3], "object_id": r[4], "timestamp": r[5].isoformat(), "curriculum_topic": r[6]}
        for r in rows
    ]
    return jsonify(statements)


@app.route("/xapi/analytics", methods=["GET"])
def get_xapi_analytics():
    """Return aggregated learner analytics."""
    conn = get_db()
    if not conn:
        return jsonify({"students": [], "struggling_concepts": [], "modules": []})
    cur = conn.cursor()

    # Per-student summary
    cur.execute("""
        SELECT actor_name, actor_email,
            COUNT(*) FILTER (WHERE verb IN ('completed', 'passed')) as mastered,
            COUNT(*) FILTER (WHERE verb = 'struggled') as struggling,
            COUNT(*) FILTER (WHERE verb = 'experienced') as viewed,
            MAX(timestamp) as last_seen
        FROM xapi_statements
        GROUP BY actor_name, actor_email
        ORDER BY last_seen DESC
    """)
    students = [
        {"name": r[0], "email": r[1], "mastered": r[2], "struggling": r[3], "viewed": r[4], "last_seen": r[5].isoformat()}
        for r in cur.fetchall()
    ]

    # Concepts with highest struggle rate
    cur.execute("""
        SELECT object_name, COUNT(*) as struggle_count
        FROM xapi_statements
        WHERE verb = 'struggled'
        GROUP BY object_name
        ORDER BY struggle_count DESC
        LIMIT 5
    """)
    struggling_concepts = [{"concept": r[0], "count": r[1]} for r in cur.fetchall()]

    # Module completion rates
    cur.execute("""
        SELECT object_id,
            COUNT(DISTINCT actor_email) FILTER (WHERE verb IN ('completed', 'passed')) as completed,
            COUNT(DISTINCT actor_email) as total_interacted
        FROM xapi_statements
        WHERE object_id LIKE 'module/%' AND object_id NOT LIKE 'module/%/%'
        GROUP BY object_id
        ORDER BY object_id
    """)
    modules = [{"module_id": r[0], "completed": r[1], "total": r[2]} for r in cur.fetchall()]

    conn.close()
    return jsonify({"students": students, "struggling_concepts": struggling_concepts, "modules": modules})


@app.route("/api/syllabus/import", methods=["POST"])
def import_syllabus():
    file = request.files.get("file")
    if not file:
        return jsonify({"error": "No file provided"}), 400

    filename = file.filename or ""
    ext = os.path.splitext(filename)[1].lower()
    if ext not in (".pdf", ".docx"):
        return jsonify({"error": "Invalid file type. Only PDF and DOCX allowed."}), 400

    file.seek(0, 2)  # seek to end
    size = file.tell()
    file.seek(0)     # reset
    if size > 10 * 1024 * 1024:
        return jsonify({"error": "File too large. Maximum size is 10MB."}), 400

    content = file.read()

    if ext == ".pdf":
        doc = fitz.open(stream=content, filetype="pdf")
        text = "\n".join(page.get_text() for page in doc)
    elif ext == ".docx":
        doc = _docx_lib.Document(io.BytesIO(content))
        text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())

    # Truncate to ~6000 chars to keep token cost low
    text = text[:6000]

    prompt = f"""You are an academic curriculum analyst. Extract structured course information from the following syllabus text.

Return ONLY valid JSON with these fields (use null if a field cannot be found):
{{
  "topic": "course name/title",
  "course_code": "e.g. CALL 301",
  "level": one of ["undergraduate-year-1","undergraduate-year-2","undergraduate-year-3","undergraduate-year-4","graduate","phd","professional"] or null,
  "audience": "discipline/field e.g. 'Applied Linguistics' or 'Business Administration'",
  "module_count": number of weeks/modules as integer or null,
  "references": [
    {{"title": "...", "url": null, "type": "academic|video|news", "reading_type": "required"}}
  ]
}}

Mark ALL extracted references as reading_type "required" — the professor chose them, so they are required.
Do not include any text outside the JSON object.

Syllabus text:
{text}"""

    try:
        if AI_PROVIDER == "gemini":
            model = genai.GenerativeModel("gemini-2.0-flash-lite")
            response = model.generate_content(prompt)
            raw = response.text
        else:
            response = openai_client.chat.completions.create(
                model="gpt-4o-mini",
                messages=[{"role": "user", "content": prompt}],
            )
            raw = response.choices[0].message.content or ""

        clean = raw.replace("```json", "").replace("```", "").strip()
        first = clean.index("{")
        last = clean.rindex("}")
        parsed = json.loads(clean[first:last + 1])
        return jsonify(parsed)
    except Exception as e:
        print(f"Syllabus import error: {e}")
        return jsonify({"error": f"Parsing failed: {str(e)}"}), 500


@app.route("/api/curriculum/export/docx", methods=["POST"])
def export_docx():
    from docx import Document as DocxDocument
    from docx.shared import Pt

    data = request.get_json()
    doc = DocxDocument()

    citation_format = data.get("citation_format", "apa")

    def fmt_citation(title, url, fmt):
        from datetime import date
        today = date.today().strftime("%Y-%m-%d")
        try:
            from urllib.parse import urlparse
            domain = urlparse(url).netloc.replace("www.", "")
        except:
            domain = url
        if fmt == "apa":
            return f"{title}. Retrieved from {url}"
        elif fmt == "mla":
            return f'"{title}." {domain}, {url}'
        elif fmt == "chicago":
            return f'"{title}." Accessed {today}. {url}'
        return title

    LEVEL_LABELS = {
        'undergraduate-year-1': 'Undergraduate — Year 1',
        'undergraduate-year-2': 'Undergraduate — Year 2',
        'undergraduate-year-3': 'Undergraduate — Year 3',
        'undergraduate-year-4': 'Undergraduate — Year 4',
        'graduate': 'Graduate',
        'phd': 'PhD',
        'professional': 'Professional',
    }
    raw_level = data.get("level", "")
    level_label = LEVEL_LABELS.get(raw_level, raw_level)

    # Title
    topic = data.get("topic", "Curriculum")
    title_para = doc.add_heading(topic, level=1)
    title_para.runs[0].bold = True

    # Course code / level / audience
    meta_parts = [p for p in [data.get("course_code"), level_label, data.get("audience")] if p]
    if meta_parts:
        doc.add_paragraph(" / ".join(meta_parts))

    # Course narrative
    course_narrative = data.get("course_narrative", "")
    if course_narrative:
        doc.add_heading("Course Narrative", level=2)
        doc.add_paragraph(course_narrative)

    # Modules
    all_readings = []
    seen_urls = set()

    for mod in data.get("modules", []):
        mod_num = mod.get("module_number", "")
        mod_title = mod.get("title", "")
        heading_text = f"Module {mod_num}: {mod_title}" if mod_num else mod_title
        doc.add_heading(heading_text, level=2)

        objectives = mod.get("learning_objectives", [])
        if objectives:
            doc.add_paragraph("Learning Objectives", style="Normal").runs[0].bold = True if doc.paragraphs[-1].runs else None
            p = doc.paragraphs[-1]
            if p.runs:
                p.runs[0].bold = True
            for obj in objectives:
                doc.add_paragraph(obj, style="List Bullet")

        readings = mod.get("recommended_readings", [])
        if readings:
            rp = doc.add_paragraph("Readings")
            if rp.runs:
                rp.runs[0].bold = True
            for r in readings:
                doc.add_paragraph(r.get('title', ''), style="List Bullet")
                key = r.get('url') or r.get('title', '')
                if key and key not in seen_urls:
                    seen_urls.add(key)
                    all_readings.append(r)

        assignments = mod.get("assignments", [])
        if assignments:
            ap = doc.add_paragraph("Assessment")
            if ap.runs:
                ap.runs[0].bold = True
            for a in assignments:
                a_name = a.get("title") or a.get("name", "")
                a_desc = a.get("task_description", "")
                name_para = doc.add_paragraph(a_name)
                if name_para.runs:
                    name_para.runs[0].bold = True
                if a_desc:
                    doc.add_paragraph(a_desc)

    # References section
    if all_readings:
        doc.add_heading("References", level=2)
        for i, r in enumerate(all_readings, start=1):
            text = fmt_citation(r.get('title', ''), r.get('url', ''), citation_format)
            doc.add_paragraph(f"[{i}] {text}")

    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)

    filename = topic.lower().replace(" ", "_") + "_curriculum.docx"
    return send_file(
        buf,
        mimetype="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        as_attachment=True,
        download_name=filename,
    )


def _slug(text: str) -> str:
    """Convert a subject name to a filesystem-safe slug (e.g. 'CALL 201' → 'call-201')."""
    return re.sub(r'[^a-z0-9]+', '-', text.strip().lower()).strip('-')


def _extract_text_from_bytes(filename: str, content: bytes) -> str:
    """Extract plain text from PDF, PPTX, or DOCX bytes."""
    ext = os.path.splitext(filename.lower())[1]
    if ext == ".pdf":
        doc = fitz.open(stream=content, filetype="pdf")
        pages = [page.get_text() for page in doc]
        doc.close()
        return "\n".join(pages)
    elif ext == ".pptx":
        from pptx import Presentation
        import io as _io
        prs = Presentation(_io.BytesIO(content))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    texts.append(shape.text)
        return "\n".join(texts)
    elif ext == ".docx":
        doc = _docx_lib.Document(io.BytesIO(content))
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    return ""


@app.route("/api/materials/ingest", methods=["POST"])
def start_ingest():
    """Accept uploaded files + subject name, kick off LightRAG ingestion in the background.

    Form fields:
        files[]  — one or more PDF/PPTX/DOCX file uploads
        subject  — string subject name (e.g. "Business Law")

    Returns immediately with {"job_id": "...", "status": "running"}.
    """
    if "files[]" not in request.files:
        return jsonify({"error": "No files[] field in request"}), 400

    subject = request.form.get("subject", "").strip()
    if not subject:
        return jsonify({"error": "Missing subject field"}), 400

    uploaded = request.files.getlist("files[]")
    if not uploaded:
        return jsonify({"error": "No files uploaded"}), 400

    # Read file bytes now (before the request context closes)
    file_data = []
    for f in uploaded:
        name = f.filename or "unknown"
        ext = os.path.splitext(name.lower())[1]
        if ext not in (".pdf", ".pptx", ".docx"):
            continue
        file_data.append((name, f.read()))

    if not file_data:
        return jsonify({"error": "No valid PDF/PPTX/DOCX files found"}), 400

    job_id = str(uuid.uuid4())
    _ingest_jobs[job_id] = {"status": "running", "progress": "Starting…", "message": ""}

    # Determine storage directory for this subject
    subject_slug = _slug(subject)
    # Map known canonical slugs to existing storage dirs
    if subject_slug == "business-law":
        storage_subdir = "lightrag_storage"
    elif subject_slug == "call":
        storage_subdir = "lightrag_storage_call"
    else:
        storage_subdir = f"lightrag_storage_{subject_slug}"

    backend_dir = os.path.dirname(os.path.abspath(__file__))
    storage_dir = os.path.normpath(os.path.join(backend_dir, "..", "data", storage_subdir))
    os.makedirs(storage_dir, exist_ok=True)

    def _run_ingest():
        total = len(file_data)
        try:
            async def _do_ingest():
                rag = _get_lightrag_instance(storage_dir)
                if storage_dir not in _initialized_instances:
                    await rag.initialize_storages()
                    _initialized_instances.add(storage_dir)
                for i, (fname, content) in enumerate(file_data, start=1):
                    _ingest_jobs[job_id]["progress"] = f"Ingesting file {i}/{total}: {fname}"
                    text = _extract_text_from_bytes(fname, content)
                    if not text.strip():
                        continue
                    # LightRAG deduplicates by hashing the raw content string.  If the
                    # same bytes were uploaded in a previous (failed) attempt, the hash
                    # is already in doc_status and the file is silently skipped.  We
                    # prepend a deterministic header containing the subject slug and
                    # filename so that:
                    #   • each file in each subject gets a unique content hash, and
                    #   • re-uploading the exact same file to the same subject is still
                    #     correctly de-duplicated (same header → same hash).
                    tagged_text = (
                        f"[source: {subject_slug} / {fname}]\n\n{text}"
                    )
                    try:
                        # lightrag-hku >=1.4 accepts an `ids` list for explicit IDs.
                        doc_id = f"doc-{subject_slug}__{os.path.splitext(fname)[0]}"
                        await rag.ainsert(tagged_text, ids=[doc_id])
                    except TypeError:
                        # Older releases don't support the ids kwarg — fall back to
                        # plain insert (the tagged_text prefix still ensures uniqueness).
                        await rag.ainsert(tagged_text)

            future = asyncio.run_coroutine_threadsafe(_do_ingest(), _bg_loop)
            future.result(timeout=600)
            _ingest_jobs[job_id]["status"] = "done"
            _ingest_jobs[job_id]["progress"] = f"Done — {total} file(s) ingested."
        except Exception as exc:
            print(f"Ingest job {job_id} failed: {exc}")
            _ingest_jobs[job_id]["status"] = "error"
            _ingest_jobs[job_id]["message"] = str(exc)

    import threading as _t
    _t.Thread(target=_run_ingest, daemon=True).start()

    return jsonify({"job_id": job_id, "status": "running"})


@app.route("/api/materials/ingest/status/<job_id>", methods=["GET"])
def ingest_status(job_id):
    """Return the current status of an ingestion job."""
    job = _ingest_jobs.get(job_id)
    if job is None:
        return jsonify({"status": "not_found"}), 404
    if job["status"] == "running":
        return jsonify({"status": "running", "progress": job.get("progress", "")})
    if job["status"] == "done":
        return jsonify({"status": "done"})
    return jsonify({"status": "error", "message": job.get("message", "Unknown error")})


init_db()
seed_mock_xapi()

if __name__ == "__main__":
    port = int(os.environ.get("PORT", 5000))
    app.run(host="0.0.0.0", port=port, debug=True)
